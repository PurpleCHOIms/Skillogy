"""Skillogy 발표자료 생성 — 3분 피치, 7장 구조."""

from pathlib import Path

from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

KOREAN_FONT = "맑은 고딕"


def set_font(run, name=KOREAN_FONT):
    """Set both Latin and East Asian font on a run so Korean renders correctly."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)

# ────────────────────────────────────────────────────────────────────────────
# Theme
# ────────────────────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x14, 0x2A)       # deep navy
BG_LIGHT = RGBColor(0xF7, 0xF8, 0xFC)      # near white
INK = RGBColor(0x10, 0x14, 0x2C)           # body text on light
INK_SOFT = RGBColor(0x4A, 0x52, 0x6B)      # secondary text
ACCENT = RGBColor(0x6E, 0x5B, 0xF6)        # purple (graph)
ACCENT_WARM = RGBColor(0xFF, 0x6B, 0x4A)   # warm coral (key numbers)
EDGE = RGBColor(0xC9, 0xCE, 0xDB)          # box border
SOFT_FILL = RGBColor(0xEE, 0xEC, 0xFF)     # tint of accent

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ────────────────────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────────────────────
def add_blank_slide(prs, bg=BG_LIGHT):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return slide


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KOREAN_FONT,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        set_font(run, font)
    return tb


def add_box(slide, x, y, w, h, *, fill=None, line=EDGE, line_w=1.0, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(line_w)
    box.shadow.inherit = False
    if radius:
        box.adjustments[0] = 0.12
    return box


def add_node(slide, x, y, w, h, label, *, fill=ACCENT, text_color=RGBColor(0xFF, 0xFF, 0xFF), size=18):
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    node.line.fill.background()
    node.shadow.inherit = False
    node.adjustments[0] = 0.5
    tf = node.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    set_font(run)
    return node


def add_edge(slide, x1, y1, x2, y2, label=None, *, color=INK_SOFT, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if label:
        cx, cy = (x1 + x2) // 2, (y1 + y2) // 2
        tb = slide.shapes.add_textbox(cx - Inches(0.7), cy - Inches(0.18), Inches(1.4), Inches(0.36))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Inches(0.04)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = color
        set_font(run, "Consolas")


def add_section_tag(slide, idx_label, name):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
             idx_label, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.65), Inches(8), Inches(0.6),
             name, size=28, bold=True, color=INK)


# ────────────────────────────────────────────────────────────────────────────
# Slide 1 — 인트로 (Skillogy 정의 + 메타포)
# ────────────────────────────────────────────────────────────────────────────
def slide_intro(prs):
    s = add_blank_slide(prs, bg=BG_DARK)

    # 좌측 상단 작은 라벨
    add_text(s, Inches(0.7), Inches(0.5), Inches(5), Inches(0.4),
             "cmux × AIM Intelligence Hackathon · 2026-04-26",
             size=12, color=RGBColor(0xA0, 0xA8, 0xC0))

    # 타이틀
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.6),
             "Skillogy", size=92, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # 한 줄 정의
    add_text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(0.6),
             "Skill + Ontology Graph — Claude Code skill 라우터",
             size=22, color=ACCENT_WARM)

    # 메타포
    add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(2.0),
             "수백 개 skill을 흩어진 카드 더미 대신,\n의미의 지도로 다시 그렸습니다.",
             size=32, color=RGBColor(0xE6, 0xE9, 0xF5))


# ────────────────────────────────────────────────────────────────────────────
# Slide 2 — 문제정의 (개인 통증 + 객관 데이터)
# ────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    s = add_blank_slide(prs)
    add_section_tag(s, "01 · PROBLEM", "Skill 트리거는 매번 lottery 다")

    # 상단: 개인 경험
    top_y = Inches(1.7)
    add_text(s, Inches(0.7), top_y, Inches(12), Inches(0.5),
             "내 skill 들이 세션마다 사라진다",
             size=24, bold=True, color=ACCENT)

    bullets = [
        "·  ~/.claude/skills 에 정성껏 만들어 둔 skill 도 세션마다 결과가 다르다",
        "·  좋은 skill 받아 깔아도 그 세션에서만 작동",
        "·  내 통증인 줄 알았는데 — 측정된 사실이었다",
    ]
    for i, b in enumerate(bullets):
        add_text(s, Inches(1.0), top_y + Inches(0.6 + 0.45 * i),
                 Inches(11.5), Inches(0.4), b, size=18, color=INK)

    # 구분선
    sep_y = Inches(4.2)
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(0.7), sep_y, Inches(12.6), sep_y)
    line.line.color.rgb = EDGE
    line.line.width = Pt(0.75)

    # 하단: 측정된 천장
    add_text(s, Inches(0.7), sep_y + Inches(0.2), Inches(12), Inches(0.5),
             "Anthropic 자체 측정 — 천장은 분명하다",
             size=24, bold=True, color=ACCENT)

    # 두 수치 카드
    card_y = sep_y + Inches(0.85)
    card_w = Inches(5.8)
    card_h = Inches(1.95)

    # Card 1
    c1 = add_box(s, Inches(0.7), card_y, card_w, card_h,
                 fill=BG_LIGHT, line=EDGE, line_w=1.5, radius=True)
    add_text(s, Inches(0.9), card_y + Inches(0.18), card_w, Inches(0.4),
             "Opus 4 · Native tool 사용 정확도", size=14, color=INK_SOFT)
    add_text(s, Inches(0.9), card_y + Inches(0.55), card_w, Inches(1.1),
             "49%", size=72, bold=True, color=ACCENT_WARM)

    # Card 2
    c2 = add_box(s, Inches(7.0), card_y, card_w, card_h,
                 fill=BG_LIGHT, line=EDGE, line_w=1.5, radius=True)
    add_text(s, Inches(7.2), card_y + Inches(0.18), card_w, Inches(0.4),
             "+ Tool Search Tool (retrieval)", size=14, color=INK_SOFT)
    add_text(s, Inches(7.2), card_y + Inches(0.55), card_w, Inches(1.1),
             "74%", size=72, bold=True, color=ACCENT_WARM)
    add_text(s, Inches(7.2), card_y + Inches(1.55), card_w, Inches(0.35),
             "← Anthropic이 retrieval 까지 동원해 도달한 천장",
             size=12, color=INK_SOFT)


# ────────────────────────────────────────────────────────────────────────────
# Slide 3 — 방법론 1: 온톨로지
# ────────────────────────────────────────────────────────────────────────────
def slide_method_ontology(prs):
    s = add_blank_slide(prs)
    add_section_tag(s, "02 · METHOD ①",
                    "Skill 시스템을 지식 그래프로 풀어냈다")

    # 캡션 우측 상단
    add_text(s, Inches(7.5), Inches(0.65), Inches(5.3), Inches(0.6),
             "Skill 자체가 1급 노드.\n문제·의도·신호 위에 정박.",
             size=14, color=INK_SOFT, align=PP_ALIGN.RIGHT)

    # 노드 좌표 (중심 기준)
    cx, cy = Inches(6.7), Inches(4.4)
    node_w, node_h = Inches(2.4), Inches(0.9)

    def place(x, y, label, fill=ACCENT):
        return add_node(s, x - node_w // 2, y - node_h // 2,
                        node_w, node_h, label, fill=fill, size=17)

    # Skill (중앙)
    skill = place(cx, cy, "Skill")

    # ProblemClass (오른쪽)
    pc_x, pc_y = cx + Inches(3.3), cy
    pc = place(pc_x, pc_y, "ProblemClass", fill=RGBColor(0x2E, 0x86, 0xC1))

    # Intent (위)
    in_x, in_y = cx, cy - Inches(2.1)
    intent = place(in_x, in_y, "Intent", fill=RGBColor(0x4A, 0x90, 0xE2))

    # Signal (아래)
    sg_x, sg_y = cx, cy + Inches(2.1)
    signal = place(sg_x, sg_y, "Signal", fill=RGBColor(0x4A, 0x90, 0xE2))

    # Skill (조합) 좌측
    sk2_x, sk2_y = cx - Inches(3.3), cy
    skill2 = place(sk2_x, sk2_y, "Skill (조합)", fill=ACCENT)

    # 엣지
    add_edge(s, cx + node_w // 2, cy, pc_x - node_w // 2, pc_y,
             label="solves", color=ACCENT_WARM, width=2.0)
    add_edge(s, in_x, in_y + node_h // 2, cx, cy - node_h // 2,
             label="demands", color=INK_SOFT)
    add_edge(s, sg_x, sg_y - node_h // 2, cx, cy + node_h // 2,
             label="triggered_by", color=INK_SOFT)
    add_edge(s, cx - node_w // 2, cy, sk2_x + node_w // 2, sk2_y,
             label="composes_with", color=INK_SOFT)

    # 범례
    leg_y = Inches(6.6)
    add_text(s, Inches(0.7), leg_y, Inches(12), Inches(0.4),
             "노드 4종 · 엣지 4종 — 라우팅·조합·안전성·관찰가능성을 한 백본에서.",
             size=14, color=INK_SOFT)


# ────────────────────────────────────────────────────────────────────────────
# Slide 4 — 방법론 2: 빌드 파이프라인
# ────────────────────────────────────────────────────────────────────────────
def slide_method_pipeline(prs):
    s = add_blank_slide(prs)
    add_section_tag(s, "02 · METHOD ②",
                    "기존 SKILL 생태계 그대로 — 우리가 추가한 건 그래프뿐")

    # 3단 파이프라인
    box_w = Inches(3.8)
    box_h = Inches(2.4)
    gap = Inches(0.55)
    total_w = box_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) // 2
    box_y = Inches(2.6)

    def step(idx, x, title, lines, fill, dark=False):
        b = add_box(s, x, box_y, box_w, box_h, fill=fill,
                    line=EDGE, line_w=1.0, radius=True)
        text_color = RGBColor(0xFF, 0xFF, 0xFF) if dark else INK
        sub_color = RGBColor(0xE6, 0xE9, 0xF5) if dark else INK_SOFT
        add_text(s, x + Inches(0.25), box_y + Inches(0.22),
                 box_w - Inches(0.5), Inches(0.4),
                 f"STEP {idx}", size=11, bold=True, color=ACCENT_WARM if not dark else ACCENT_WARM)
        add_text(s, x + Inches(0.25), box_y + Inches(0.55),
                 box_w - Inches(0.5), Inches(0.6),
                 title, size=20, bold=True, color=text_color)
        for i, ln in enumerate(lines):
            add_text(s, x + Inches(0.25),
                     box_y + Inches(1.25 + 0.36 * i),
                     box_w - Inches(0.5), Inches(0.36),
                     ln, size=13, color=sub_color)

    # Step 1
    x1 = start_x
    step(1, x1, "SKILL.md 스캔",
         ["~/.claude/skills/**/SKILL.md",
          "공식 frontmatter + body",
          "Anthropic 표준 그대로"],
         fill=BG_LIGHT)

    # Step 2
    x2 = x1 + box_w + gap
    step(2, x2, "LLM 추출",
         ["Haiku 4.5 가 컨텍스트로 받음",
          "엔티티(노드) + 릴레이션(엣지)",
          "스키마 자동 추출"],
         fill=SOFT_FILL)

    # Step 3
    x3 = x2 + box_w + gap
    step(3, x3, "Neo4j Graph DB",
         ["Skillogy KG 구축",
          "Cypher 로 traversal",
          "FastAPI · MCP 로 노출"],
         fill=ACCENT, dark=True)

    # 화살표
    arr_y = box_y + box_h // 2
    for ax_start in [x1 + box_w, x2 + box_w]:
        arr = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     ax_start + Inches(0.05), arr_y,
                                     ax_start + gap - Inches(0.05), arr_y)
        arr.line.color.rgb = ACCENT
        arr.line.width = Pt(2.5)

    # 캡션
    add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
             "기존 SKILL 표준을 깨지 않는다. Skillogy 는 그 위에 그래프 레이어를 더할 뿐.",
             size=16, color=INK_SOFT, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Slide 5 — 벤치마크 (이미지 placeholder)
# ────────────────────────────────────────────────────────────────────────────
def slide_benchmark(prs):
    s = add_blank_slide(prs)
    add_section_tag(s, "03 · BENCHMARK",
                    "Trigger Accuracy — Native vs Skillogy")

    # placeholder box
    p_x = Inches(2.0)
    p_y = Inches(1.8)
    p_w = Inches(9.3)
    p_h = Inches(4.7)
    add_box(s, p_x, p_y, p_w, p_h,
            fill=BG_LIGHT, line=EDGE, line_w=1.5, radius=True)

    add_text(s, p_x, p_y + Inches(1.6), p_w, Inches(0.6),
             "[ 차트 이미지 삽입 영역 ]",
             size=24, bold=True, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_text(s, p_x, p_y + Inches(2.4), p_w, Inches(0.5),
             "Native vs Skillogy · Trigger Accuracy 2-bar",
             size=16, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_text(s, p_x, p_y + Inches(3.0), p_w, Inches(0.4),
             "발표 직전 차트를 본 위치에 PNG 로 삽입",
             size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
             "동일 SKILL 풀에서 Native(=Anthropic progressive disclosure) vs Skillogy 적용 후 비교.",
             size=13, color=INK_SOFT)


# ────────────────────────────────────────────────────────────────────────────
# Slide 6 — 라이브 데모 안내
# ────────────────────────────────────────────────────────────────────────────
def slide_demo(prs):
    s = add_blank_slide(prs, bg=BG_DARK)

    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "04 · LIVE", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
             "Live Demo", size=36, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # 두 분할 박스
    half_w = Inches(5.9)
    half_h = Inches(3.2)
    half_y = Inches(2.4)
    gap = Inches(0.4)
    left_x = (SLIDE_W - half_w * 2 - gap) // 2
    right_x = left_x + half_w + gap

    # Native
    nb = add_box(s, left_x, half_y, half_w, half_h,
                 fill=RGBColor(0x1B, 0x22, 0x40), line=EDGE, line_w=0.75, radius=True)
    add_text(s, left_x + Inches(0.3), half_y + Inches(0.25),
             half_w - Inches(0.6), Inches(0.4),
             "WITHOUT SKILLOGY", size=12, bold=True,
             color=RGBColor(0xFF, 0x9F, 0x80))
    add_text(s, left_x + Inches(0.3), half_y + Inches(0.65),
             half_w - Inches(0.6), Inches(0.7),
             "Native ❌", size=32, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, left_x + Inches(0.3), half_y + Inches(1.55),
             half_w - Inches(0.6), Inches(1.5),
             "동일 발화 →\n잘못된 / 빈 skill 트리거",
             size=18, color=RGBColor(0xC0, 0xC8, 0xE0))

    # Skillogy
    sb = add_box(s, right_x, half_y, half_w, half_h,
                 fill=ACCENT, line=EDGE, line_w=0.75, radius=True)
    add_text(s, right_x + Inches(0.3), half_y + Inches(0.25),
             half_w - Inches(0.6), Inches(0.4),
             "WITH SKILLOGY", size=12, bold=True,
             color=RGBColor(0xFF, 0xE2, 0xC8))
    add_text(s, right_x + Inches(0.3), half_y + Inches(0.65),
             half_w - Inches(0.6), Inches(0.7),
             "Skillogy ✅", size=32, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, right_x + Inches(0.3), half_y + Inches(1.55),
             half_w - Inches(0.6), Inches(1.5),
             "그래프 traversal →\n맞는 skill 트리거",
             size=18, color=RGBColor(0xF1, 0xEB, 0xFF))

    # 진행 안내
    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
             "터미널 2분할 (15s)  →  Web UI 그래프 시각화 (20s)",
             size=18, color=RGBColor(0xE6, 0xE9, 0xF5),
             align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Slide 7 — 클로징
# ────────────────────────────────────────────────────────────────────────────
def slide_closing(prs):
    s = add_blank_slide(prs, bg=BG_DARK)

    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.0),
             "Skill 트리거의 천장은",
             size=44, color=RGBColor(0xC0, 0xC8, 0xE0),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.0),
             "LLM 머릿속에 있었다.",
             size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(1.0),
             "Skillogy 는 그걸 그래프로 끌어냈다.",
             size=48, bold=True, color=ACCENT_WARM,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "Skillogy · Skill + Ontology Graph",
             size=14, color=RGBColor(0x9A, 0xA2, 0xBE),
             align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Build
# ────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_intro(prs)
    slide_problem(prs)
    slide_method_ontology(prs)
    slide_method_pipeline(prs)
    slide_benchmark(prs)
    slide_demo(prs)
    slide_closing(prs)

    out = Path(__file__).parent / "skillogy.pptx"
    prs.save(out)
    print(f"saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
